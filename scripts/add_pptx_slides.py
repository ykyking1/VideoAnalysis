"""experiment_results pptx dosyasında slayt 1 ve 2'yi (1-indexed) olduğu gibi
bırakır; slayt 3 ve 4'ü, mevcut slayt 0 (tablo+3 tek satır stat callout
şablonu) ve slayt 3'ün (karşılaştırma tablosu şablonu) XML'i kopyalanarak
üretilen yeni içerikle değiştirir - toplam slayt sayısı 4'te sabit kalır,
format/font/renk aynen korunur.

Kullanım: python scripts/add_pptx_slides.py
"""
import copy

from pptx import Presentation

PPTX_PATH = r"C:\Users\PC_4150_YD26\Downloads\experiment_results_2026-07-24_amt_07_58_11.pptx"
BACKUP_PATH = r"C:\Users\PC_4150_YD26\Downloads\experiment_results_2026-07-24_amt_07_58_11.backup.pptx"


def duplicate_slide(prs, source_index):
    source = prs.slides[source_index]
    new_slide = prs.slides.add_slide(source.slide_layout)
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in source.shapes:
        new_slide.shapes._spTree.append(copy.deepcopy(shp._element))
    return new_slide


def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rId = slides[index].rId
    prs.part.drop_rel(rId)
    xml_slides.remove(slides[index])


def set_run(shape, para_idx, run_idx, text):
    shape.text_frame.paragraphs[para_idx].runs[run_idx].text = text


def set_cell(table, row, col, text):
    table.cell(row, col).text_frame.paragraphs[0].runs[0].text = text


def build_speed_accuracy_slide(prs):
    slide = duplicate_slide(prs, 0)
    shapes = list(slide.shapes)
    # shapes: 0=top bar, 1=title, 2=subtitle, 3=table, 4/5 6/7 8/9 = 3 stat callouts, 10=bottom insight

    set_run(shapes[1], 0, 0, "Hybrid Query Pipeline \u2014 Measured Speed & Accuracy")
    set_run(
        shapes[2], 0, 0,
        "Separate 6-video / 95-chunk corpus \u2014 own dataset, not Study 1's DiDeMo set. "
        "Self-retrieval proxy metric (N=8), not directly comparable to Study 1/2's 68-query "
        "ground-truth eval. Ollama query parsing (CPU) + ClickHouse hybrid search, GT1030 4GB.",
    )

    table = shapes[3].table
    rows = [
        ("Metric", "X-CLIP (512d)", "VideoCLIP-XL (768d)"),
        ("LLM query parsing (Ollama, CPU)", "13.8 s", "13.8 s"),
        ("Query embedding (warm)", "80 ms", "229 ms"),
        ("ClickHouse hybrid search (95 chunks)", "57 ms", "57 ms"),
        ("Interval merge", "<1 ms", "<1 ms"),
        ("Total query latency (warm)", "~13.9 s", "~14.1 s"),
        ("Storage / chunk (measured)", "2125 B", "3124 B"),
        ("Recall@1 (self-retrieval, N=8)", "12.5%", "25.0%"),
        ("Recall@5 (self-retrieval, N=8)", "62.5%", "75.0%"),
        ("Mean rank (of 95)", "9.2", "4.8"),
    ]
    for r, (a, b, c) in enumerate(rows):
        set_cell(table, r, 0, a)
        set_cell(table, r, 1, b)
        set_cell(table, r, 2, c)

    set_run(shapes[4], 0, 0, "57 ms")
    set_run(shapes[5], 0, 0, "ClickHouse hybrid query (filter + vector, single call) \u2014 negligible vs. the LLM step")
    set_run(shapes[6], 0, 0, "99%")
    set_run(shapes[7], 0, 0, "of total query latency is CPU-bound LLM parsing \u2014 not the hybrid search itself")
    set_run(shapes[8], 0, 0, "+12.5 pts")
    set_run(shapes[9], 0, 0, "Recall@1 gain with VideoCLIP-XL over X-CLIP, same hybrid pipeline")

    set_run(
        shapes[10], 0, 0,
        "The hybrid query itself is fast and effectively free (~57 ms regardless of filter count); "
        "today's bottleneck is CPU-bound LLM parsing (~13.8 s), a driver/hardware limitation of this "
        "test rig (Ollama CUDA build needs driver 570+, host has 560.94) \u2014 not a property of the hybrid design.",
    )
    return slide


def build_hybrid_vs_standard_slide(prs):
    slide = duplicate_slide(prs, 3)
    shapes = list(slide.shapes)
    # shapes: 0=title, 1=subtitle, 2=table, 3/4 5/6 7/8 = 3 stat callouts (2-line captions), 9=bottom insight

    set_run(shapes[0], 0, 0, "Hybrid vs. Standard Architecture \u2014 Structural Advantages")
    set_run(
        shapes[1], 0, 0,
        "Standard = single vector database, semantic-only. Hybrid = structured filter + vector "
        "similarity in one ClickHouse query (this project's design).",
    )

    table = shapes[2].table
    rows = [
        ("Capability", "Standard (vector-only)", "Hybrid (this project)"),
        ("Deterministic constraints (speed, altitude, time)", "Left to embedding \u2014 probabilistic", "ClickHouse WHERE \u2014 exact match"),
        ("Filter + vector in one query", "Separate DB, pre/post-filter dilemma", "Same row, single ClickHouse query"),
        ("Concept not in schema", "\u2014", "Falls back to semantic-only, no hard failure"),
        ("Result shape", "Raw top-k chunks", "Interval-merged, continuous time ranges"),
    ]
    for r, (a, b, c) in enumerate(rows):
        set_cell(table, r, 0, a)
        set_cell(table, r, 1, b)
        set_cell(table, r, 2, c)

    set_run(shapes[3], 0, 0, "0")
    set_run(shapes[4], 0, 0, "Separate vector-DB joins needed")
    set_run(shapes[4], 1, 0, "Filter + vector combined in one ClickHouse row \u2014 no pre/post-filter dilemma")

    set_run(shapes[5], 0, 0, "57 ms")
    set_run(shapes[6], 0, 0, "Measured \u2014 this session")
    set_run(shapes[6], 1, 0, "Combined filter+vector ClickHouse query, 95 chunks \u2014 no added cost from filters")

    set_run(shapes[7], 0, 0, "Deterministic")
    set_run(shapes[8], 0, 0, "Design property, not benchmarked")
    set_run(shapes[8], 1, 0, "Structured filters return exact matches, not probabilistic rankings")

    set_run(
        shapes[9], 0, 0,
        "Real telemetry data doesn't exist yet in this corpus, so the accuracy contribution of "
        "structured filtering is unmeasured (TELEMETRY_FILTERS_ENABLED=False) \u2014 but its architectural "
        "cost is negligible: the same ClickHouse query runs at the same speed whether filters are active or not.",
    )
    return slide


def main():
    # Orijinal 4 slaytlı halden (yedek) başla - slayt 0,1 (Embedding/Chunking
    # Comparison) dokunulmadan kalacak, slayt 2,3 (Why Hybrid / Cost Profile)
    # yeni içerikle değiştirilecek.
    prs = Presentation(BACKUP_PATH)
    assert len(prs.slides) == 4, f"beklenmedik slayt sayisi: {len(prs.slides)}"

    build_speed_accuracy_slide(prs)       # slayt 0'ı klonlayıp sona ekler (index 4)
    build_hybrid_vs_standard_slide(prs)   # slayt 3'ü klonlayıp sona ekler (index 5)

    # Eski slayt 3 ve 2'yi sil (yüksek index'ten başla) - kalan: [0,1,4,5] -> [0,1,2,3]
    delete_slide(prs, 3)
    delete_slide(prs, 2)

    prs.save(PPTX_PATH)
    print(f"Kaydedildi: {PPTX_PATH} ({len(prs.slides)} slayt)")


if __name__ == "__main__":
    main()
